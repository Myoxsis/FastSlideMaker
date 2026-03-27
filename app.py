from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.responses import HTMLResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

from layout_engine import layout_for_slide
from llm_client import LLMConfig, MockLLMClient, OllamaClient
from project_store import list_projects, load_project, save_project
from prompt_parser import build_system_prompt, build_user_prompt
from slide_planner import normalize_deck, select_template_for_slide

app = FastAPI(title="Fast Slide Maker", version="0.2.0")
app.mount("/static", StaticFiles(directory="static"), name="static")

INDEX_HTML = Path("templates/index.html").read_text(encoding="utf-8")
SETTINGS_PATH = Path("samples/llm_settings.json")


class LLMRuntimeConfig(BaseModel):
    base_url: str = "http://localhost:11434"
    endpoint: str = "/api/chat"
    model: str = "llama3"
    temperature: float = 0.2
    max_tokens: int = 1800
    timeout_seconds: float = 120.0
    use_mock: bool = False


class GenerateRequest(LLMRuntimeConfig):
    prompt: str
    slide_count: int = 6
    reference_context: str = ""


class RegenerateSlideRequest(LLMRuntimeConfig):
    prompt: str
    deck: dict[str, Any]
    slide_id: str
    max_tokens: int = 800
    reference_context: str = ""


class ChatRequest(LLMRuntimeConfig):
    message: str
    reference_context: str = ""


@app.get("/", response_class=HTMLResponse)
def index() -> str:
    return INDEX_HTML


def _get_llm(cfg: LLMRuntimeConfig):
    if cfg.use_mock:
        return MockLLMClient(), "mock"

    client = OllamaClient(
        LLMConfig(
            base_url=cfg.base_url,
            endpoint=cfg.endpoint,
            model=cfg.model,
            temperature=cfg.temperature,
            max_tokens=cfg.max_tokens,
            timeout_seconds=cfg.timeout_seconds,
        )
    )
    ok, reason = client.health_check()
    if not ok:
        return MockLLMClient(), f"fallback:{reason}"
    return client, "ollama"


def _default_settings() -> dict[str, Any]:
    return LLMRuntimeConfig().model_dump()


def _append_reference_context(prompt: str, reference_context: str) -> str:
    context = reference_context.strip()
    if not context:
        return prompt
    return (
        f"{prompt}\n\n"
        "Reference context from an existing PowerPoint deck. "
        "Use this only as optional context, prioritize the user request:\n"
        f"{context}"
    )


@app.get("/api/settings")
def get_settings():
    if not SETTINGS_PATH.exists():
        return _default_settings()
    return json.loads(SETTINGS_PATH.read_text(encoding="utf-8"))


@app.post("/api/settings")
def save_settings(config: LLMRuntimeConfig):
    SETTINGS_PATH.parent.mkdir(parents=True, exist_ok=True)
    SETTINGS_PATH.write_text(json.dumps(config.model_dump(), indent=2), encoding="utf-8")
    return {"saved": str(SETTINGS_PATH), "settings": config.model_dump()}


@app.post("/api/generate")
def generate(req: GenerateRequest):
    llm, mode = _get_llm(req)
    system_prompt = build_system_prompt()
    user_prompt = build_user_prompt(req.prompt, req.slide_count)
    user_prompt = _append_reference_context(user_prompt, req.reference_context)

    try:
        candidate = llm.generate_structured_json(system_prompt, user_prompt)
        deck = normalize_deck(candidate)
    except Exception as exc:  # noqa: BLE001
        if isinstance(exc, TimeoutError):
            detail = (
                "Failed to generate deck: request timed out. "
                "Increase timeout/model speed or reduce slide count/max tokens."
            )
        else:
            detail = f"Failed to generate deck: {exc}"
        raise HTTPException(status_code=500, detail=detail) from exc

    payload = deck.model_dump(by_alias=True)
    for slide in payload["slides"]:
        slide["template"] = select_template_for_slide(slide["slide_type"])
        slide["render_layout"] = layout_for_slide(slide)

    return {"mode": mode, "deck": payload}


@app.post("/api/regenerate-slide")
def regenerate_slide(req: RegenerateSlideRequest):
    llm, mode = _get_llm(req)

    target = next((s for s in req.deck.get("slides", []) if s.get("id") == req.slide_id), None)
    if not target:
        raise HTTPException(status_code=404, detail="slide_id not found")

    context = json.dumps({"presentation": req.deck.get("presentation"), "target_slide": target}, indent=2)
    system_prompt = build_system_prompt()
    user_prompt = (
        build_user_prompt(req.prompt, 1)
        + "\nRegenerate ONLY one slide with matching id and compatible narrative context.\n"
        + context
    )
    user_prompt = _append_reference_context(user_prompt, req.reference_context)

    try:
        candidate = llm.generate_structured_json(system_prompt, user_prompt)
        regen_deck = normalize_deck(candidate).model_dump(by_alias=True)
        replacement = regen_deck["slides"][0]
        replacement["id"] = req.slide_id
        replacement["template"] = select_template_for_slide(replacement["slide_type"])
        replacement["render_layout"] = layout_for_slide(replacement)
    except Exception as exc:  # noqa: BLE001
        raise HTTPException(status_code=500, detail=f"Failed to regenerate slide: {exc}") from exc

    return {"mode": mode, "slide": replacement}


@app.post("/api/chat")
def chat(req: ChatRequest):
    llm, mode = _get_llm(req)
    try:
        message = _append_reference_context(req.message, req.reference_context)
        reply = llm.chat(message)
    except Exception as exc:  # noqa: BLE001
        raise HTTPException(status_code=500, detail=f"Failed to chat: {exc}") from exc
    return {"mode": mode, "reply": reply}


@app.post("/api/context/powerpoint")
async def extract_powerpoint_context(file: UploadFile = File(...)):
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported")

    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise HTTPException(status_code=500, detail="python-pptx is not installed") from exc

    raw = await file.read()
    if not raw:
        raise HTTPException(status_code=400, detail="Uploaded file is empty")

    from io import BytesIO

    try:
        presentation = Presentation(BytesIO(raw))
    except Exception as exc:  # noqa: BLE001
        raise HTTPException(status_code=400, detail=f"Failed to parse PowerPoint file: {exc}") from exc

    slide_summaries: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            cleaned = text.strip()
            if cleaned:
                text_parts.append(cleaned)
        if text_parts:
            combined = " | ".join(text_parts)
            slide_summaries.append(f"Slide {index}: {combined}")

    if not slide_summaries:
        raise HTTPException(status_code=400, detail="No extractable text found in PowerPoint")

    joined = "\n".join(slide_summaries)
    if len(joined) > 8000:
        joined = f"{joined[:8000]}\n... (truncated)"

    return {
        "filename": file.filename,
        "slide_count": len(presentation.slides),
        "slides_with_text": len(slide_summaries),
        "reference_context": joined,
    }


@app.post("/api/save/{project_name}")
def save(project_name: str, payload: dict[str, Any]):
    path = save_project(payload, project_name)
    return {"saved": str(path)}


@app.get("/api/load/{project_name}")
def load(project_name: str):
    try:
        return load_project(project_name)
    except FileNotFoundError as exc:
        raise HTTPException(status_code=404, detail="project not found") from exc


@app.get("/api/projects")
def projects():
    return {"projects": list_projects()}
