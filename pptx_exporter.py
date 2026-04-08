from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from schema import Deck


class PptxExporter:
    def __init__(self) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def export(self, deck: Deck, output_path: Path) -> Path:
        for slide in deck.slides:
            ppt_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            self._add_title(ppt_slide, slide.title, slide.subtitle)
            if slide.slide_type == "process_flow":
                self._draw_process_flow(ppt_slide, slide)
            elif slide.slide_type == "layered_architecture":
                self._draw_layered_architecture(ppt_slide, slide)
            elif slide.slide_type == "roadmap":
                self._draw_roadmap(ppt_slide, slide)
            else:
                self._draw_bullets(ppt_slide, slide)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output_path))
        return output_path

    @staticmethod
    def _add_title(slide, title: str, subtitle: str):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.6))
        tf = title_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(15, 23, 42)

        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.4))
            sub_tf = sub_box.text_frame
            sub_tf.text = subtitle
            sub_tf.paragraphs[0].font.size = Pt(14)
            sub_tf.paragraphs[0].font.color.rgb = RGBColor(71, 85, 105)

    @staticmethod
    def _draw_process_flow(slide, semantic_slide):
        left = 0.7
        top = 1.7
        width = 2.6
        height = 1.2
        gap = 0.45

        for i, step in enumerate(semantic_slide.flow_steps[:4]):
            x = Inches(left + i * (width + gap))
            rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(top), Inches(width), Inches(height))
            rect.text = f"{step.title}\n{step.detail}"
            rect.fill.solid()
            rect.fill.fore_color.rgb = RGBColor(219, 234, 254)
            rect.line.color.rgb = RGBColor(29, 78, 216)
            if i < len(semantic_slide.flow_steps[:4]) - 1:
                slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    x + Inches(width),
                    Inches(top + height / 2),
                    x + Inches(width + gap),
                    Inches(top + height / 2),
                )

    @staticmethod
    def _draw_layered_architecture(slide, semantic_slide):
        top = 1.6
        for i, layer in enumerate(semantic_slide.layers[:5]):
            y = Inches(top + i * 1.0)
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), y, Inches(11.2), Inches(0.8))
            band.fill.solid()
            band.fill.fore_color.rgb = RGBColor(224 - i * 12, 242 - i * 8, 254 - i * 6)
            band.line.color.rgb = RGBColor(14, 116, 144)
            text = f"{layer.name}: " + " | ".join(layer.components)
            band.text_frame.text = text

    @staticmethod
    def _draw_roadmap(slide, semantic_slide):
        y = Inches(3.2)
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.0), y, Inches(12.2), y)
        for i, phase in enumerate(semantic_slide.roadmap_phases[:4]):
            x = Inches(1.5 + i * 3.0)
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.9), Inches(0.6), Inches(0.6))
            marker.fill.solid()
            marker.fill.fore_color.rgb = RGBColor(2, 132, 199)
            label = slide.shapes.add_textbox(x - Inches(0.4), Inches(3.6), Inches(2.0), Inches(1.2))
            label.text_frame.text = f"{phase.name}\n{phase.timeframe}\n- " + "\n- ".join(phase.outcomes[:2])

    @staticmethod
    def _draw_bullets(slide, semantic_slide):
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12), Inches(5.5))
        tf = box.text_frame
        tf.text = ""

        bullets = semantic_slide.narrative or ["No semantic bullets provided."]
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
