"""Deterministic exporter from semantic slide JSON to native editable PPTX."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from textwrap import fill
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.models.schemas import (
    ArchitectureLayer,
    Milestone,
    MilestoneStatus,
    ProcessStep,
    RoadmapPhase,
    SemanticPresentation,
    SemanticSlide,
    SlideType,
    Swimlane,
    TextBlock,
)
from app.services.designer import design_semantic_presentation



@dataclass(frozen=True)
class Theme:
    """PowerPoint theme constants for predictable visual output."""

    bg = RGBColor(248, 251, 255)
    title = RGBColor(18, 42, 76)
    body = RGBColor(36, 52, 71)
    muted = RGBColor(90, 112, 136)
    accent = RGBColor(47, 128, 237)
    accent_2 = RGBColor(111, 66, 193)
    success = RGBColor(26, 127, 55)
    warning = RGBColor(176, 125, 0)
    danger = RGBColor(186, 36, 41)
    border = RGBColor(203, 213, 225)
    white = RGBColor(255, 255, 255)

    title_size_pt = 28
    subtitle_size_pt = 14
    heading_size_pt = 16
    body_size_pt = 12
    small_size_pt = 10


class Spacing:
    """Spacing helpers for deterministic layout math."""

    slide_margin_x = Inches(0.55)
    slide_margin_y = Inches(0.35)
    title_height = Inches(0.8)
    objective_height = Inches(0.45)
    content_top_gap = Inches(0.15)
    card_gap = Inches(0.18)


def _wrap_text(text: str, *, max_chars: int) -> str:
    """Soft-wrap text deterministically by character count."""

    cleaned = " ".join(text.split())
    if not cleaned:
        return ""
    return "\n".join(fill(line, width=max_chars) for line in cleaned.split("\n"))


class PptxExporter:
    """Export semantic slides to editable native PowerPoint primitives."""

    def __init__(self) -> None:
        self.theme = Theme()

    def export_deck(self, semantic: SemanticPresentation | dict[str, Any], output_path: str | Path) -> Path:
        """Export the entire semantic deck to a .pptx file and return the path."""

        model = semantic if isinstance(semantic, SemanticPresentation) else SemanticPresentation.model_validate(semantic)
        deck = design_semantic_presentation(model.normalized())

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        ordered = {slide.id: slide for slide in deck.slides}
        for slide_id in deck.slide_order:
            slide = ordered[slide_id]
            pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_chrome(pptx_slide, slide)
            self._render_semantic_slide(pptx_slide, slide)

        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        return path

    def _add_slide_chrome(self, pptx_slide: Any, slide: SemanticSlide) -> float:
        """Render title + optional objective + callouts and return content top coordinate."""

        left = Spacing.slide_margin_x
        width = Inches(12.2)

        title_box = pptx_slide.shapes.add_textbox(left, Spacing.slide_margin_y, width, Spacing.title_height)
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        run = title_frame.paragraphs[0].add_run()
        run.text = _wrap_text(slide.title, max_chars=58)
        run.font.bold = True
        run.font.size = Pt(self.theme.title_size_pt)
        run.font.color.rgb = self.theme.title

        current_top = Spacing.slide_margin_y + Spacing.title_height

        if slide.objective:
            objective_box = pptx_slide.shapes.add_textbox(left, current_top, width, Spacing.objective_height)
            objective_frame = objective_box.text_frame
            objective_frame.clear()
            objective_frame.word_wrap = True
            objective_run = objective_frame.paragraphs[0].add_run()
            objective_run.text = _wrap_text(slide.objective, max_chars=95)
            objective_run.font.size = Pt(self.theme.subtitle_size_pt)
            objective_run.font.color.rgb = self.theme.muted
            current_top += Spacing.objective_height

        if slide.text_blocks:
            current_top = self._add_text_labels(pptx_slide, slide.text_blocks, current_top + Spacing.content_top_gap)

        return float(current_top + Spacing.content_top_gap)

    def _add_text_labels(self, pptx_slide: Any, text_blocks: list[TextBlock], top: float) -> float:
        left = Spacing.slide_margin_x
        width = Inches(12.2)
        card_h = Inches(0.65)

        for block in text_blocks[:2]:
            card = pptx_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = Theme.white
            card.line.color.rgb = Theme.border
            frame = card.text_frame
            frame.word_wrap = True
            frame.clear()
            p = frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            text = f"{block.label}: {block.text}" if block.label else block.text
            run = p.add_run()
            run.text = _wrap_text(text, max_chars=125)
            run.font.size = Pt(self.theme.small_size_pt)
            run.font.color.rgb = self.theme.body
            top += card_h + Inches(0.08)

        return float(top)

    def _render_semantic_slide(self, pptx_slide: Any, slide: SemanticSlide) -> None:
        positions = slide.layout_hints.element_positions

        if slide.process:
            self._render_process_flow(pptx_slide, slide.process.steps, positions)
            return

        if slide.swimlanes:
            self._render_swimlanes(pptx_slide, slide.swimlanes.lanes, positions)
            return

        if slide.architecture:
            self._render_layered_architecture(pptx_slide, slide.architecture.layers, positions)
            return

        if slide.roadmap:
            self._render_roadmap(pptx_slide, slide.roadmap.phases, positions)
            return

        if slide.type == SlideType.CONTENT and slide.text_blocks:
            self._render_basic_bullets(pptx_slide, slide.text_blocks, positions)

    def _render_basic_bullets(self, pptx_slide: Any, blocks: list[TextBlock], positions: dict[str, Any]) -> None:
        fallback = {"x": 0.75, "y": 1.95, "w": 12.0, "h": 4.8}
        first_box = positions.get(blocks[0].id, fallback) if blocks else fallback
        box = pptx_slide.shapes.add_textbox(Inches(first_box["x"]), Inches(first_box["y"]), Inches(first_box["w"]), Inches(first_box["h"]))
        frame = box.text_frame
        frame.word_wrap = True
        frame.clear()

        for idx, block in enumerate(blocks):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = f"• {_wrap_text(block.text, max_chars=100)}"
            paragraph.level = 0
            paragraph.font.size = Pt(self.theme.body_size_pt)
            paragraph.font.color.rgb = self.theme.body

    def _render_process_flow(self, pptx_slide: Any, steps: list[ProcessStep], positions: dict[str, Any]) -> None:
        left = Inches(0.75)
        top = Inches(2.0)
        width = Inches(2.35)
        height = Inches(1.2)
        gap = Inches(0.34)

        step_boxes: list[Any] = []
        max_cards = min(len(steps), 4)

        if len(steps) > 4:
            width = Inches(1.95)
            gap = Inches(0.22)
            max_cards = min(len(steps), 5)

        for idx, step in enumerate(steps[:max_cards]):
            hinted = positions.get(step.id)
            x = Inches(hinted["x"]) if hinted else left + idx * (width + gap)
            y = Inches(hinted["y"]) if hinted else top
            card_w = Inches(hinted["w"]) if hinted else width
            card_h = Inches(hinted["h"]) if hinted else height
            box = pptx_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
            box.fill.solid()
            box.fill.fore_color.rgb = Theme.white
            box.line.color.rgb = Theme.accent

            tf = box.text_frame
            tf.clear()
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = f"{idx + 1}. {step.label}"
            p0.font.bold = True
            p0.font.size = Pt(self.theme.heading_size_pt)
            p0.font.color.rgb = self.theme.title

            desc = step.description or ""
            p1 = tf.add_paragraph()
            p1.text = _wrap_text(desc, max_chars=28)
            p1.font.size = Pt(self.theme.small_size_pt)
            p1.font.color.rgb = self.theme.body

            if step.owner:
                p2 = tf.add_paragraph()
                p2.text = f"Owner: {step.owner}"
                p2.font.size = Pt(self.theme.small_size_pt)
                p2.font.color.rgb = self.theme.muted

            step_boxes.append(box)

        self._connect_boxes(pptx_slide, step_boxes)

    def _connect_boxes(self, pptx_slide: Any, boxes: list[Any]) -> None:
        if len(boxes) < 2:
            return

        for left_box, right_box in zip(boxes, boxes[1:], strict=False):
            x1 = left_box.left + left_box.width
            y1 = left_box.top + left_box.height / 2
            x2 = right_box.left
            y2 = right_box.top + right_box.height / 2
            connector = pptx_slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
            connector.line.color.rgb = Theme.accent
            connector.line.width = Pt(2)
            connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID

    def _render_layered_architecture(self, pptx_slide: Any, layers: list[ArchitectureLayer], positions: dict[str, Any]) -> None:
        left = Inches(1.0)
        top = Inches(1.95)
        width = Inches(11.2)
        layer_h = Inches(1.28)

        for idx, layer in enumerate(layers[:4]):
            hinted = positions.get(layer.id)
            y = Inches(hinted["y"]) if hinted else top + idx * (layer_h + Inches(0.18))
            card_left = Inches(hinted["x"]) if hinted else left
            card_width = Inches(hinted["w"]) if hinted else width
            card_height = Inches(hinted["h"]) if hinted else layer_h
            rect = pptx_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, card_left, y, card_width, card_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = Theme.white if idx % 2 == 0 else RGBColor(241, 246, 255)
            rect.line.color.rgb = Theme.border

            tf = rect.text_frame
            tf.clear()
            tf.word_wrap = True
            header = tf.paragraphs[0]
            header.text = layer.name
            header.font.bold = True
            header.font.size = Pt(self.theme.heading_size_pt)
            header.font.color.rgb = self.theme.accent_2

            detail = tf.add_paragraph()
            detail.text = _wrap_text(layer.responsibility, max_chars=110)
            detail.font.size = Pt(self.theme.body_size_pt)
            detail.font.color.rgb = self.theme.body

            if layer.components:
                components = tf.add_paragraph()
                components.text = f"Components: {', '.join(layer.components[:6])}"
                components.font.size = Pt(self.theme.small_size_pt)
                components.font.color.rgb = self.theme.muted

    def _render_swimlanes(self, pptx_slide: Any, lanes: list[Swimlane], positions: dict[str, Any]) -> None:
        lane_count = min(max(len(lanes), 1), 6)
        left = Inches(0.75)
        top = Inches(1.95)
        width = Inches(12.0)
        gap = Inches(0.12)
        total_height = Inches(4.95)
        lane_height = (total_height - gap * (lane_count - 1)) / lane_count
        label_width = Inches(2.1)

        for idx, lane in enumerate(lanes[:lane_count]):
            hinted = positions.get(lane.id)
            lane_top = Inches(hinted["y"]) if hinted else top + idx * (lane_height + gap)
            lane_row_h = Inches(hinted["h"]) if hinted else lane_height

            label_card = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                left,
                lane_top,
                label_width,
                lane_row_h,
            )
            label_card.fill.solid()
            label_card.fill.fore_color.rgb = RGBColor(233, 242, 255) if idx % 2 == 0 else RGBColor(240, 247, 255)
            label_card.line.color.rgb = Theme.accent

            label_tf = label_card.text_frame
            label_tf.clear()
            label_tf.word_wrap = True
            label_p = label_tf.paragraphs[0]
            label_p.alignment = PP_ALIGN.CENTER
            label_p.text = _wrap_text(lane.lane_label, max_chars=20)
            label_p.font.bold = True
            label_p.font.size = Pt(self.theme.body_size_pt)
            label_p.font.color.rgb = self.theme.title

            body_card = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left + label_width,
                lane_top,
                width - label_width,
                lane_row_h,
            )
            body_card.fill.solid()
            body_card.fill.fore_color.rgb = Theme.white if idx % 2 == 0 else RGBColor(250, 252, 255)
            body_card.line.color.rgb = Theme.border

            body_tf = body_card.text_frame
            body_tf.clear()
            body_tf.word_wrap = True

            if not lane.items:
                paragraph = body_tf.paragraphs[0]
                paragraph.text = "No activities defined."
                paragraph.font.size = Pt(self.theme.small_size_pt)
                paragraph.font.color.rgb = self.theme.muted
                continue

            max_items = min(len(lane.items), 5)
            for item_idx, item in enumerate(lane.items[:max_items]):
                paragraph = body_tf.paragraphs[0] if item_idx == 0 else body_tf.add_paragraph()
                detail = f": {item.detail}" if item.detail else ""
                paragraph.text = f"• {_wrap_text(item.label + detail, max_chars=86)}"
                paragraph.level = 0
                paragraph.font.size = Pt(self.theme.small_size_pt)
                paragraph.font.color.rgb = self.theme.body

    def _render_roadmap(self, pptx_slide: Any, phases: list[RoadmapPhase], positions: dict[str, Any]) -> None:
        axis_left = Inches(1.1)
        axis_top = Inches(3.2)
        axis_width = Inches(10.8)

        axis = pptx_slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            axis_left,
            axis_top,
            axis_left + axis_width,
            axis_top,
        )
        axis.line.color.rgb = Theme.border
        axis.line.width = Pt(2.25)

        count = min(max(len(phases), 1), 5)
        spacing = axis_width / count
        marker_r = Inches(0.13)

        for idx, phase in enumerate(phases[:count]):
            hinted = positions.get(phase.id)
            center_x = axis_left + spacing * idx + spacing / 2 if not hinted else Inches(hinted["x"] + hinted["w"] / 2)
            circle = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                center_x - marker_r,
                axis_top - marker_r,
                marker_r * 2,
                marker_r * 2,
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = Theme.accent
            circle.line.color.rgb = Theme.accent

            card_top = Inches(1.92) if idx % 2 == 0 else Inches(3.45)
            card_h = Inches(1.35)
            card_w = Inches(2.2)
            card_left = center_x - Inches(1.1)
            if hinted:
                card_top = Inches(hinted["y"])
                card_h = Inches(hinted["h"])
                card_w = Inches(hinted["w"])
                card_left = Inches(hinted["x"])
            card = pptx_slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                card_left,
                card_top,
                card_w,
                card_h,
            )
            card.fill.solid()
            card.fill.fore_color.rgb = Theme.white
            card.line.color.rgb = Theme.border

            tf = card.text_frame
            tf.clear()
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = phase.name
            p0.font.bold = True
            p0.font.size = Pt(self.theme.body_size_pt)
            p0.font.color.rgb = self.theme.title

            p1 = tf.add_paragraph()
            p1.text = _wrap_text(phase.objective, max_chars=30)
            p1.font.size = Pt(self.theme.small_size_pt)
            p1.font.color.rgb = self.theme.body

            for milestone in phase.milestones[:2]:
                p = tf.add_paragraph()
                p.text = self._milestone_line(milestone)
                p.font.size = Pt(9)
                p.font.color.rgb = self._milestone_color(milestone)

            start_y = axis_top
            end_y = card_top + (card_h if idx % 2 == 0 else 0)
            link = pptx_slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, center_x, start_y, center_x, end_y)
            link.line.color.rgb = Theme.border
            link.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT

    @staticmethod
    def _milestone_line(milestone: Milestone) -> str:
        return f"• {milestone.label} ({milestone.target_period})"

    @staticmethod
    def _milestone_color(milestone: Milestone) -> RGBColor:
        if milestone.status == MilestoneStatus.COMPLETE:
            return Theme.success
        if milestone.status == MilestoneStatus.AT_RISK:
            return Theme.danger
        if milestone.status == MilestoneStatus.IN_PROGRESS:
            return Theme.warning
        return Theme.muted


def export_semantic_deck_to_pptx(
    semantic_json: SemanticPresentation | dict[str, Any],
    output_path: str | Path,
) -> Path:
    """Convenience function for full-deck export."""

    return PptxExporter().export_deck(semantic_json, output_path)
